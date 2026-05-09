from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── colour palette ──────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT   = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2  = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0xB0, 0xC4, 0xD8)
CARD_BG  = RGBColor(0x1A, 0x2E, 0x44)   # slightly lighter navy

# ── helpers ─────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_bullet_box(slide, title, bullets, l, t, w, h, title_size=13, bullet_size=11):
    """Card with title bar + bullet list."""
    # card background
    card = add_rect(slide, l, t, w, h, CARD_BG)
    # top accent bar
    add_rect(slide, l, t, w, 0.07, ACCENT)
    # title
    add_text(slide, title, l+0.15, t+0.10, w-0.3, 0.35,
             size=title_size, bold=True, color=ACCENT)
    # bullets
    tb = slide.shapes.add_textbox(Inches(l+0.18), Inches(t+0.48),
                                   Inches(w-0.36), Inches(h-0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"▸  {b}"
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = GRAY
        p.space_after = Pt(4)

def slide_title_accent(slide, title, subtitle=None):
    """Full-width title + optional subtitle at top."""
    add_rect(slide, 0, 0, 13.33, 1.1, CARD_BG)
    add_rect(slide, 0, 1.05, 13.33, 0.06, ACCENT)
    add_text(slide, title, 0.4, 0.08, 12.5, 0.65,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.68, 12.5, 0.38,
                 size=14, color=ACCENT2, align=PP_ALIGN.LEFT)

blank = prs.slide_layouts[6]   # completely blank layout

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)

# decorative rectangles
add_rect(s, 0,   0,   13.33, 0.08, ACCENT)
add_rect(s, 0,   7.42, 13.33, 0.08, ACCENT)
add_rect(s, 0,   0,   0.12, 7.5,  ACCENT)
add_rect(s, 13.21, 0, 0.12, 7.5,  ACCENT)

# large centre card
add_rect(s, 1.5, 1.6, 10.33, 4.5, CARD_BG)
add_rect(s, 1.5, 1.6, 10.33, 0.1, ACCENT)

add_text(s, "SIMD Bitonic Sort", 1.7, 1.8, 10.0, 1.1,
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "High-Performance Sorting Using SIMD Parallelism",
         1.7, 2.85, 10.0, 0.6,
         size=20, color=ACCENT2, align=PP_ALIGN.CENTER)

add_rect(s, 4.5, 3.6, 4.33, 0.04, ACCENT)

add_text(s, "Parallel & Distributed Computing", 1.7, 3.85, 10.0, 0.45,
         size=15, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "SIMD • Bitonic Sort • AVX / NEON • Merge Sort",
         1.7, 4.3, 10.0, 0.45,
         size=13, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s, "2026", 1.7, 5.3, 10.0, 0.4,
         size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — Introduction
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
slide_title_accent(s, "Introduction", "Problem Area & Why It Matters")

# left big card — problem
add_rect(s, 0.3, 1.25, 6.1, 5.9, CARD_BG)
add_rect(s, 0.3, 1.25, 6.1, 0.07, ACCENT)
add_text(s, "The Problem", 0.5, 1.35, 5.7, 0.4, size=15, bold=True, color=ACCENT)

problems = [
    ("Sorting is everywhere", "Graphics, databases, physics engines, compression — all depend on fast sorting."),
    ("std::sort is too slow", "Sequential comparison-based sort processes one element pair at a time."),
    ("Small arrays are neglected", "Existing fast sorts target large arrays. Small arrays (< 200 elements) are called millions of times per second in real workloads."),
    ("CPU potential wasted", "Modern CPUs can compare 8 floats simultaneously via AVX registers — but standard sort never uses this."),
]
y = 1.85
for title, desc in problems:
    add_text(s, f"▸  {title}", 0.5, y, 5.7, 0.28, size=12, bold=True, color=ACCENT2)
    add_text(s, f"   {desc}",  0.5, y+0.26, 5.7, 0.38, size=10, color=GRAY)
    y += 0.72

# right big card — importance
add_rect(s, 6.9, 1.25, 6.1, 5.9, CARD_BG)
add_rect(s, 6.9, 1.25, 6.1, 0.07, ACCENT)
add_text(s, "Why It Is Important", 7.1, 1.35, 5.7, 0.4, size=15, bold=True, color=ACCENT)

importance = [
    ("7× speedup over std::sort", "Benchmarked on real hardware — not theoretical."),
    ("Cache-friendly", "Sorting stays inside CPU registers — almost zero memory traffic during sort."),
    ("Branch-free hot path", "Fixed comparison network means no mispredicted branches — the CPU pipeline stays full."),
    ("Real use cases", "Image compression (8×8 blocks), KD-tree construction, physics simulations — all sort small float arrays repeatedly."),
    ("Portable", "Single header file works on x86 (AVX) and ARM (NEON) with the same logic."),
]
y = 1.85
for title, desc in importance:
    add_text(s, f"▸  {title}", 7.1, y, 5.7, 0.28, size=12, bold=True, color=ACCENT2)
    add_text(s, f"   {desc}",  7.1, y+0.26, 5.7, 0.38, size=10, color=GRAY)
    y += 0.72

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — Core Idea
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
slide_title_accent(s, "Core Idea", "Central Contribution & Parallelism Strategy")

cards = [
    ("What Kind of Parallelism?",
     ["Data parallelism — the same min/max operation applied to 8 floats simultaneously",
      "AVX register holds 8 floats; one instruction compares all 8 pairs at once",
      "NEON register holds 4 floats; same principle on ARM",
      "No task parallelism — single threaded, but maximally wide per instruction"]),
    ("Parallelization Strategy",
     ["Decomposition: array split into vectors (chunks of 8 floats each)",
      "Each vector sorted independently with simd_sort_1V (no communication needed)",
      "Orchestration: sorted vectors merged using bitonic merge network",
      "For large arrays: tiled merge sort — each tile ≤ 192 floats sorted by SIMD, then merged"]),
    ("Data Locality & Cache Use",
     ["Small sort (≤192 floats) fits entirely in L1 cache — zero cache misses during sort",
      "Load once from RAM → sort inside registers → store once back to RAM",
      "Tile size (192) chosen so working set stays in L1/L2 cache",
      "Merge phase uses temporary buffers — only necessary memory copies made"]),
    ("Bitonic Network — The Core",
     ["Fixed sequence of compare-and-swap operations — no data-dependent branching",
      "Round 1: compare distance-1 pairs → Round 2: distance-3 → ... → Round 6: full sort",
      "Every round: shuffle → parallel min → parallel max → blend result",
      "Same network topology regardless of input values — predictable pipeline behaviour"]),
]

positions = [(0.3, 1.25), (6.72, 1.25), (0.3, 4.4), (6.72, 4.4)]
for (l, t), (title, bullets) in zip(positions, cards):
    add_bullet_box(s, title, bullets, l, t, 6.2, 2.9, title_size=13, bullet_size=10)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Approach
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
slide_title_accent(s, "Approach", "How the Problem Is Solved Step by Step")

# timeline-style vertical flow on left
steps = [
    ("STEP 1", "Load",    "Read floats from RAM into SIMD registers (8 at a time). Partial vectors padded with +∞ so they sort to the end and are discarded on store."),
    ("STEP 2", "Sort 1V", "simd_sort_1V: 6 rounds of shuffle → min/max → blend inside one register. Sorts 8 floats with ~18 SIMD instructions. Zero memory access."),
    ("STEP 3", "Sort NV", "simd_sort_NV builds recursively: sort_4V = sort_2V(a,b) + sort_2V(c,d) + merge_4V. Scales from 2 to 24 vectors (up to 192 floats)."),
    ("STEP 4", "Merge",   "simd_merge_2V_sorted: reverse one vector, min/max against the other → bitonic sequence → aftermerge rounds finish it. All in registers."),
    ("STEP 5", "Store",   "Write sorted registers back to RAM. Partial last vector writes only the real elements — +∞ padding discarded."),
]

y = 1.3
for tag, name, desc in steps:
    add_rect(s, 0.3, y, 1.0, 0.72, ACCENT)
    add_text(s, tag,  0.3, y+0.04, 1.0, 0.28, size=8,  bold=True, color=BG,   align=PP_ALIGN.CENTER)
    add_text(s, name, 0.3, y+0.3,  1.0, 0.3,  size=11, bold=True, color=BG,   align=PP_ALIGN.CENTER)
    add_rect(s, 1.3, y+0.18, 5.5, 0.45, CARD_BG)
    add_text(s, desc, 1.45, y+0.13, 5.2, 0.55, size=9.5, color=GRAY)
    if y < 4.8:
        add_rect(s, 0.72, y+0.72, 0.16, 0.22, ACCENT)
    y += 0.94

# right side — code-style diagram
add_rect(s, 7.1, 1.25, 5.9, 5.9, CARD_BG)
add_rect(s, 7.1, 1.25, 5.9, 0.07, ACCENT)
add_text(s, "Call Hierarchy", 7.3, 1.35, 5.5, 0.38, size=13, bold=True, color=ACCENT)

hierarchy = """simd_merge_sort(array, n)
│
├─ n ≤ 192
│    └─ simd_small_sort()
│         ├─ load into registers
│         ├─ simd_sort_NV()
│         │    ├─ simd_sort_1V()   [6 rounds]
│         │    ├─ simd_sort_1V()
│         │    └─ simd_merge_NV_sorted()
│         └─ store to memory
│
└─ n > 192
     └─ merge_sort()
          ├─ sort left  (recurse / simd_small_sort)
          ├─ sort right (recurse / simd_small_sort)
          └─ merge_arrays()  [SIMD merge]"""

tb = s.shapes.add_textbox(Inches(7.25), Inches(1.78), Inches(5.6), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.text = hierarchy
p.font.size = Pt(9.5)
p.font.color.rgb = ACCENT2
p.font.name = "Courier New"

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Gap / Opportunity
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
slide_title_accent(s, "Gap / Opportunity", "Limitations & Further Optimization Potential")

gaps = [
    ("1  Single-Threaded Only",
     "The entire sort runs on one CPU core. All other cores sit idle. For large arrays this is the biggest missed opportunity — the tile sort phase is embarrassingly parallel.",
     "HIGH"),
    ("2  Branchy SIMD Merge",
     "merge_arrays() contains an if/else that decides whether to load from left or right half each iteration. This breaks CPU branch prediction and stalls the pipeline unnecessarily.",
     "HIGH"),
    ("3  No Prefetching",
     "When processing large arrays the CPU frequently stalls waiting for data from RAM. _mm_prefetch hints could load the next chunk into cache while sorting the current one — essentially free latency hiding.",
     "MED"),
    ("4  AVX-512 Not Used",
     "Modern Intel server and desktop CPUs support AVX-512 (16 floats per register vs 8). Porting the bitonic network to AVX-512 would double throughput per instruction with identical algorithm logic.",
     "MED"),
    ("5  Hardcoded Tile Size",
     "simd_small_sort_max() always returns 192 regardless of CPU. Optimal tile size should match L1 cache size — a smaller L1 cache will suffer unnecessary cache misses with the current fixed value.",
     "LOW"),
]

impact_colors = {"HIGH": RGBColor(0xFF,0x6B,0x6B), "MED": RGBColor(0xFF,0xD1,0x66), "LOW": RGBColor(0x6B,0xFF,0xB8)}

y = 1.28
for title, desc, impact in gaps:
    add_rect(s, 0.3, y, 12.5, 1.02, CARD_BG)
    add_rect(s, 0.3, y, 0.07, 1.02, impact_colors[impact])
    # impact badge
    add_rect(s, 11.7, y+0.28, 1.0, 0.38, impact_colors[impact])
    add_text(s, impact, 11.7, y+0.3, 1.0, 0.3, size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_text(s, title, 0.5, y+0.06, 10.8, 0.32, size=12, bold=True, color=WHITE)
    add_text(s, desc,  0.5, y+0.38, 11.0, 0.55, size=9.5, color=GRAY)
    y += 1.14

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — Proposed Approach
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
slide_title_accent(s, "Proposed Approach", "What We Will Implement, Adapt & Evaluate")

# 3-column layout
cols = [
    ("Implement\nOpenMP Threading",
     ["Parallelize the tile-sort phase using #pragma omp parallel for",
      "Each thread independently calls simd_small_sort on its tile",
      "Tiles are fully independent — zero synchronization needed during sort phase",
      "Merge phase remains sequential as baseline",
      "Expected: near-linear speedup with core count on large arrays"]),
    ("Implement\nBranchless Merge",
     ["Replace the if/else branch in merge_arrays with masked SIMD loads",
      "Use _mm256_blendv_ps to select from left or right without branching",
      "Eliminates branch misprediction stalls in the merge hot loop",
      "Keeps the rest of the SIMD merge logic identical",
      "Expected: measurable throughput improvement on large arrays"]),
    ("Evaluate\nAVX-512 Width",
     ["Port simd_sort_1V to AVX-512 — 16 floats per register instead of 8",
      "Bitonic network rounds stay identical, just wider shuffle/blend masks",
      "simd_small_sort_max grows from 192 to 384 floats",
      "Benchmark on AVX-512 capable hardware vs current AVX baseline",
      "Expected: up to 2x throughput on supported CPUs"]),
]

x = 0.3
for title, bullets in cols:
    add_rect(s, x, 1.25, 4.15, 5.95, CARD_BG)
    add_rect(s, x, 1.25, 4.15, 0.07, ACCENT)
    add_text(s, title, x+0.18, 1.35, 3.8, 0.65, size=13, bold=True, color=ACCENT)
    tb = s.shapes.add_textbox(Inches(x+0.2), Inches(2.1), Inches(3.75), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"▸  {b}"
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.space_after = Pt(6)
    x += 4.45

# bottom bar — what we won't do
add_rect(s, 0.3, 7.08, 12.73, 0.32, RGBColor(0x1A,0x1A,0x2E))
add_text(s, "Out of scope:  MPI (network overhead kills small-array benefit)  •  CUDA (copy latency dominates for small arrays)  •  Key-value sorting",
         0.5, 7.1, 12.3, 0.28, size=9, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Preliminary Plan
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
slide_title_accent(s, "Preliminary Plan", "Tools, Languages, Platforms & Implementation Roadmap")

# left column — tools & platform
add_rect(s, 0.3, 1.25, 4.5, 5.95, CARD_BG)
add_rect(s, 0.3, 1.25, 4.5, 0.07, ACCENT)
add_text(s, "Tools & Platform", 0.5, 1.35, 4.1, 0.38, size=13, bold=True, color=ACCENT)

tools = [
    ("Language",  "C99 / C++17"),
    ("SIMD",      "Intel AVX2 intrinsics (_mm256_*)  +  ARM NEON (arm_neon.h)"),
    ("Threading", "OpenMP 4.5  (#pragma omp parallel for)"),
    ("Compiler",  "GCC / Clang  with  -mavx2 -fopenmp -O3"),
    ("Platform",  "x86-64 Linux / Windows  (AVX2 required)"),
    ("Benchmark", "sokol_time.h  (already in repo)  +  custom harness"),
    ("Baseline",  "std::sort  vs  original simd_merge_sort  vs  new version"),
    ("Profiling", "perf stat  /  Intel VTune  for cache miss & IPC analysis"),
]

y = 1.85
for label, val in tools:
    add_text(s, label, 0.5,  y, 1.3,  0.28, size=10, bold=True, color=ACCENT2)
    add_text(s, val,   1.85, y, 2.8,  0.28, size=10, color=GRAY)
    add_rect(s, 0.5, y+0.3, 3.9, 0.01, RGBColor(0x25,0x3A,0x52))
    y += 0.62

# right column — timeline
add_rect(s, 5.2, 1.25, 7.9, 5.95, CARD_BG)
add_rect(s, 5.2, 1.25, 7.9, 0.07, ACCENT)
add_text(s, "Implementation Roadmap", 5.4, 1.35, 7.5, 0.38, size=13, bold=True, color=ACCENT)

phases = [
    ("Week 1", "Baseline & Benchmarking",
     "Set up benchmark harness. Profile existing code with perf/VTune. Record throughput for array sizes 8 → 10M floats. Establish baseline numbers for std::sort and current simd_merge_sort."),
    ("Week 2", "OpenMP Tile Parallelism",
     "Add #pragma omp parallel for over tile sort phase. Verify correctness with randomized test suite. Benchmark across 1, 2, 4, 8, 16 threads. Measure scaling efficiency."),
    ("Week 3", "Branchless SIMD Merge",
     "Replace if/else in merge_arrays with masked SIMD loads (_mm256_blendv_ps). Benchmark merge-dominated workloads (large arrays already sorted in tiles). Measure branch-miss reduction via perf stat."),
    ("Week 4", "AVX-512 Port + Final Evaluation",
     "Port simd_sort_1V and bitonic network to __m512 (16-wide). Run full benchmark suite. Compile results table: std::sort vs original vs +OpenMP vs +branchless vs +AVX-512."),
]

y = 1.85
for week, phase, desc in phases:
    add_rect(s, 5.35, y, 1.15, 0.88, ACCENT)
    add_text(s, week,  5.35, y+0.04, 1.15, 0.3,  size=9,  bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_text(s, phase, 5.35, y+0.34, 1.15, 0.45, size=7.5,bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_rect(s, 6.5, y+0.15, 6.3, 0.65, RGBColor(0x12,0x24,0x38))
    add_text(s, desc, 6.6, y+0.1, 6.1, 0.75, size=9, color=GRAY)
    if y < 4.6:
        add_rect(s, 5.85, y+0.88, 0.15, 0.18, ACCENT)
    y += 1.08

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Summary / Thank You
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_rect(s, 0, 0, 13.33, 0.08, ACCENT)
add_rect(s, 0, 7.42, 13.33, 0.08, ACCENT)
add_rect(s, 0, 0, 0.12, 7.5, ACCENT)
add_rect(s, 13.21, 0, 0.12, 7.5, ACCENT)

add_text(s, "Summary", 0.5, 0.6, 12.5, 0.75, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, 4.5, 1.38, 4.33, 0.06, ACCENT)

summary_points = [
    ("Problem",   "Small float arrays are sorted millions of times per second in real workloads — std::sort wastes CPU potential"),
    ("Core Idea", "Bitonic sort mapped onto SIMD registers — 8 comparisons per instruction, all work in registers, near-zero memory traffic"),
    ("Approach",  "Fixed comparison network (shuffle → min → max → blend) × 6 rounds per vector, scaled up to 24 vectors via recursive merge"),
    ("Gap",       "Single-threaded, branchy merge, no prefetch, no AVX-512, hardcoded tile size"),
    ("Proposed",  "OpenMP tile parallelism + branchless merge + AVX-512 port — targeting measurable speedup over the already-fast baseline"),
]

y = 1.6
for label, text in summary_points:
    add_rect(s, 1.2, y, 1.5, 0.6, ACCENT)
    add_text(s, label, 1.2, y+0.08, 1.5, 0.42, size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_rect(s, 2.7, y+0.06, 9.7, 0.48, CARD_BG)
    add_text(s, text, 2.85, y+0.06, 9.4, 0.48, size=10.5, color=GRAY)
    y += 0.78

add_text(s, "Thank You", 0.5, 6.7, 12.5, 0.55, size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ── save ────────────────────────────────────────────────────────
out = r"d:\PDC\simd_bitonic-master\SIMD_Bitonic_Sort_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
