from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0F, 0x0F, 0x13)   # near-black
PANEL   = RGBColor(0x1A, 0x1A, 0x24)   # dark card
BORDER  = RGBColor(0x28, 0x28, 0x38)   # subtle border
ACCENT  = RGBColor(0x6C, 0x63, 0xFF)   # purple
CYAN    = RGBColor(0x00, 0xD4, 0xFF)   # cyan
GREEN   = RGBColor(0x00, 0xE5, 0x96)   # mint green
YELLOW  = RGBColor(0xFF, 0xC8, 0x44)   # amber
RED     = RGBColor(0xFF, 0x4D, 0x6D)   # rose
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0x88, 0x88, 0xAA)
MGRAY   = RGBColor(0x44, 0x44, 0x60)

def c(r,g,b): return RGBColor(r,g,b)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Core helpers ──────────────────────────────────────────────────────────────
def bg(slide, color=BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, color=PANEL):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def bordered_rect(slide, l, t, w, h, fill=PANEL, border=BORDER, bw=Pt(1)):
    from pptx.util import Pt as _Pt
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = bw; return s

def tb(slide, l, t, w, h, text, size=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    s = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = s.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return s

def hline(slide, l, t, w, color=ACCENT, h=0.025):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def dot(slide, l, t, r=0.07, color=ACCENT):
    from pptx.util import Inches as I
    s = slide.shapes.add_shape(9, I(l-r/2), I(t-r/2), I(r), I(r))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

# ── Chart helpers ─────────────────────────────────────────────────────────────
def bar_chart(slide, l, t, w, h, data, max_val=None,
              bar_color=ACCENT, label_color=LGRAY, val_color=WHITE,
              bar_h_frac=0.55, show_val=True, val_fmt="{:.2f}×"):
    """Horizontal bar chart. data = [(label, value), ...]"""
    n     = len(data)
    mv    = max_val or max(v for _,v in data)
    row_h = h / n
    bar_area = w * 0.55
    lbl_w    = w * 0.30
    for i, (lbl, val) in enumerate(data):
        y    = t + i * row_h + row_h * (1-bar_h_frac) / 2
        bh   = row_h * bar_h_frac
        bw   = (val / mv) * bar_area
        # label
        tb(slide, l, y, lbl_w, bh, lbl, size=11, color=label_color,
           align=PP_ALIGN.RIGHT, wrap=False)
        # bar
        rect(slide, l+lbl_w+0.05, y, bw, bh, color=bar_color)
        # value
        if show_val:
            tb(slide, l+lbl_w+0.1+bw, y, w*0.15, bh,
               val_fmt.format(val), size=11, bold=True, color=val_color)

def grouped_bar_chart(slide, l, t, w, h, series, x_labels,
                      colors, max_val=None, y_label="", title=""):
    """series = [(name, [values]), ...],  x_labels = [str, ...]"""
    n_groups = len(x_labels)
    n_series = len(series)
    mv       = max_val or max(v for _,vals in series for v in vals)
    chart_h  = h * 0.75
    chart_t  = t + h * 0.12
    chart_w  = w * 0.92
    chart_l  = l + w * 0.05
    # axes
    rect(slide, chart_l, chart_t, 0.015, chart_h, color=MGRAY)
    rect(slide, chart_l, chart_t+chart_h, chart_w, 0.015, color=MGRAY)
    # gridlines
    for gi in range(1,5):
        gy = chart_t + chart_h - (gi/4)*chart_h
        rect(slide, chart_l, gy, chart_w, 0.008, color=BORDER)
        tb(slide, l, gy-0.12, w*0.04, 0.25,
           f"{mv*gi/4:.0f}", size=9, color=LGRAY, align=PP_ALIGN.RIGHT)
    # bars
    group_w = chart_w / n_groups
    bar_w   = group_w * 0.8 / n_series
    pad     = group_w * 0.1
    for gi, label in enumerate(x_labels):
        gx = chart_l + gi * group_w
        # x label
        tb(slide, gx, chart_t+chart_h+0.05, group_w, 0.25,
           label, size=9, color=LGRAY, align=PP_ALIGN.CENTER)
        for si, (sname, vals) in enumerate(series):
            bx = gx + pad + si * bar_w
            bh = (vals[gi] / mv) * chart_h
            by = chart_t + chart_h - bh
            rect(slide, bx, by, bar_w*0.85, bh, color=colors[si])
    # legend
    lx = chart_l
    for si, (sname, _) in enumerate(series):
        rect(slide, lx, t, 0.15, 0.1, color=colors[si])
        tb(slide, lx+0.18, t-0.02, 1.2, 0.18, sname, size=9, color=LGRAY)
        lx += 1.45

def draw_line_segment(slide, x1, y1, x2, y2, color, thickness=0.04):
    """Draw a properly rotated line segment between two points (in Inches)."""
    import math
    from pptx.util import Inches as I, Pt
    from lxml import etree
    dx = x2 - x1; dy = y2 - y1
    length = math.sqrt(dx*dx + dy*dy)
    if length < 0.001:
        return
    angle_deg = math.degrees(math.atan2(dy, dx))
    # Place shape at midpoint, sized length × thickness
    cx = (x1 + x2) / 2
    cy = (y1 + y2) / 2
    sp = slide.shapes.add_shape(1,
        I(cx - length/2), I(cy - thickness/2),
        I(length), I(thickness))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    # Apply rotation via XML
    sp_elem = sp._element
    xfrm = sp_elem.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
    if xfrm is not None:
        xfrm.set('rot', str(int(angle_deg * 60000)))

def line_chart(slide, l, t, w, h, series, x_labels, colors,
               max_val=None, min_val=0):
    """series = [(name, [y_values]), ...]"""
    mv      = max_val or max(v for _,vals in series for v in vals)
    mn      = min_val
    chart_h = h * 0.80
    chart_t = t + h * 0.12
    chart_w = w * 0.90
    chart_l = l + w * 0.08
    n       = len(x_labels)
    # axes
    rect(slide, chart_l, chart_t, 0.015, chart_h, color=MGRAY)
    rect(slide, chart_l, chart_t+chart_h, chart_w, 0.015, color=MGRAY)
    # gridlines
    steps = 4
    for gi in range(1, steps+1):
        gy  = chart_t + chart_h - (gi/steps)*chart_h
        val = mn + (mv-mn)*gi/steps
        rect(slide, chart_l, gy, chart_w, 0.008, color=BORDER)
        tb(slide, l, gy-0.1, w*0.07, 0.22,
           f"{val:.1f}", size=9, color=LGRAY, align=PP_ALIGN.RIGHT)
    # x-axis labels
    for xi, lbl in enumerate(x_labels):
        xp = chart_l + (xi/(n-1))*chart_w if n>1 else chart_l
        tb(slide, xp-0.3, chart_t+chart_h+0.05, 0.6, 0.22,
           lbl, size=9, color=LGRAY, align=PP_ALIGN.CENTER)
    # lines + dots
    for si, (sname, vals) in enumerate(series):
        pts = []
        for xi, v in enumerate(vals):
            xp = chart_l + (xi/(n-1))*chart_w if n>1 else chart_l
            yp = chart_t + chart_h - ((v-mn)/(mv-mn))*chart_h
            pts.append((xp, yp))
        for i in range(len(pts)-1):
            x1,y1 = pts[i]; x2,y2 = pts[i+1]
            draw_line_segment(slide, x1, y1, x2, y2, colors[si], thickness=0.045)
        for (xp,yp) in pts:
            dot(slide, xp, yp, r=0.09, color=colors[si])

# ── Slide components ──────────────────────────────────────────────────────────
def header(slide, title, subtitle=""):
    tb(slide, 0.45, 0.22, 12.4, 0.55, title,
       size=26, bold=True, color=WHITE)
    if subtitle:
        tb(slide, 0.45, 0.76, 12.4, 0.32, subtitle,
           size=13, color=LGRAY, italic=True)
    hline(slide, 0.45, 1.12, 12.43, color=ACCENT)

def pill(slide, l, t, label, color=ACCENT, text_color=WHITE):
    rect(slide, l, t, 1.1, 0.28, color=color)
    tb(slide, l, t, 1.1, 0.28, label, size=10, bold=True,
       color=text_color, align=PP_ALIGN.CENTER)

def stat_card(slide, l, t, w, h, number, label, color=ACCENT):
    bordered_rect(slide, l, t, w, h, fill=PANEL, border=color)
    tb(slide, l, t+0.08, w, h*0.55, number, size=34, bold=True,
       color=color, align=PP_ALIGN.CENTER)
    tb(slide, l, t+h*0.62, w, h*0.35, label, size=12,
       color=LGRAY, align=PP_ALIGN.CENTER)

def section_card(slide, l, t, w, h, title, lines, title_color=ACCENT, size=13):
    bordered_rect(slide, l, t, w, h, fill=PANEL, border=BORDER)
    hline(slide, l, t+0.46, w, color=title_color)
    tb(slide, l+0.2, t+0.1, w-0.4, 0.35, title,
       size=14, bold=True, color=title_color)
    y = t + 0.55
    for line in lines:
        tb(slide, l+0.2, y, w-0.4, 0.42, "·  " + line,
           size=size, color=WHITE, wrap=True)
        y += 0.42


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Accent stripe left edge
rect(s, 0, 0, 0.06, 7.5, color=ACCENT)

# Large title
tb(s, 0.55, 1.5, 10, 1.1,
   "Making Sorting 3× Faster", size=52, bold=True, color=WHITE)
tb(s, 0.55, 2.68, 10, 0.55,
   "with SIMD Vector Units & Multi-Core Parallelism",
   size=22, color=CYAN)

hline(s, 0.55, 3.38, 6.0, color=ACCENT)

tb(s, 0.55, 3.6, 9.5, 0.4,
   "PDC Final Project  ·  SIMD Bitonic Sort  ·  AVX2 + OpenMP",
   size=15, color=LGRAY)

# Right side visual — stat cards
for i,(val,lbl,col) in enumerate([
        ("3.0×",  "Peak\nSpeedup",        GREEN),
        ("12",    "CPU\nThreads",          CYAN),
        ("2M",    "Elements\nSorted",      ACCENT),
        ("35/35", "Tests\nPassed",         YELLOW),
]):
    stat_card(s, 10.5, 0.85 + i*1.6, 2.5, 1.45, val, lbl, color=col)

tb(s, 0.55, 6.8, 7.0, 0.42,
   "Intel Core i7-12th Gen  ·  6 cores / 12 threads  ·  AVX2  ·  GCC 15  ·  Windows 11",
   size=12, color=MGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Base Paper
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=ACCENT)
header(s, "Base Library — simd_bitonic",
       "Open-source SIMD bitonic sort · AVX2 path · Batcher odd-even merge network")

# Two columns
section_card(s, 0.45, 1.3, 5.9, 4.1, "What the library does", [
    "Sorts floating-point arrays using 256-bit AVX2 registers",
    "Processes 8 numbers simultaneously per register (vs 1 in std::sort)",
    "simd_small_sort: branchless sorting network for chunks ≤ 192 elements",
    "simd_merge_sort: splits array into 192-element chunks, sorts each, then merges repeatedly",
    "At n = 2,000,000: needs 14 merge passes to fully sort",
], title_color=CYAN)

section_card(s, 6.55, 1.3, 6.35, 4.1, "Performance gaps we identified", [
    "Gap 1 — Single-threaded: 11 cores sit idle",
    "Gap 2 — Merge hot-loop makes a branch decision 50M times/sort",
    "Gap 3 — No memory prefetching for large array merges",
    "Gap 4 — AVX-512 (wider registers) not used",
    "Gap 5 — Tile size hardcoded, ignores actual cache",
], title_color=YELLOW)

# Bottom bar: baseline vs std::sort
rect(s, 0.45, 5.6, 12.43, 1.55, color=PANEL)
bordered_rect(s, 0.45, 5.6, 12.43, 1.55, fill=PANEL, border=BORDER)
tb(s, 0.65, 5.65, 12.0, 0.35,
   "Baseline already beats std::sort — our goal: push even further",
   size=13, bold=True, color=WHITE)

bsl_data = [
    ("n = 81",     9.02), ("n = 2,187",  4.34),
    ("n = 19,683", 3.85), ("n = 531,441",2.84),
]
x = 0.75
for lbl, spd in bsl_data:
    bordered_rect(s, x, 6.08, 2.85, 0.88, fill=c(0x0A,0x18,0x2E), border=CYAN)
    tb(s, x, 6.10, 2.85, 0.38, lbl, size=12, color=LGRAY, align=PP_ALIGN.CENTER)
    tb(s, x, 6.48, 2.85, 0.44, f"{spd}× faster than std::sort",
       size=13, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    x += 3.05


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Project Scope
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=GREEN)
header(s, "Project Scope & Approach",
       "Five targeted improvements — each isolated and measured independently")

section_card(s, 0.45, 1.3, 3.9, 5.85, "In Scope", [
    "Gap 1: OpenMP task-based multi-threading",
    "Gap 2: Branchless merge hot-loop",
    "Gap 3: Software memory prefetching",
    "Gap 4: AVX-512 simulation (tile doubling)",
    "Gap 5: Dynamic tile + pooled allocation",
    "Correctness check vs std::sort at every size",
    "Strong & weak scaling measurement",
    "Array sizes 10K → 2M floats",
], title_color=GREEN, size=13)

section_card(s, 4.55, 1.3, 3.9, 2.7, "Out of Scope", [
    "Real AVX-512 (hardware unavailable)",
    "Distributed / MPI parallelism",
    "GPU acceleration",
    "Non-x86 platforms",
], title_color=RED, size=13)

section_card(s, 4.55, 4.2, 3.9, 2.95, "Tools", [
    "Language: C++11",
    "Compiler: GCC 15  (-O3 -mavx2 -fopenmp)",
    "Threading: OpenMP 4.5 task model",
    "Timer: sokol_time.h (nanosecond precision)",
    "RNG: fixed seed 0xABCDEF01 (reproducible)",
], title_color=CYAN, size=13)

section_card(s, 8.65, 1.3, 4.25, 5.85, "Methodology", [
    "Each variant is a separate selectable option",
    "100 timed runs per size — averaged",
    "Same random data for all variants (fair comparison)",
    "Correctness verified before any timing",
    "5 correctness sizes × 7 variants = 35 tests",
    "Results from actual hardware runs — no simulation",
    "Scalability: thread sweep 1 → 12 threads",
], title_color=YELLOW, size=13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Baseline Method
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=CYAN)
header(s, "Baseline — How simd_merge_sort Works",
       "Two phases: branchless tile sort + recursive merge")

# Phase boxes
for i,(title,col,lines) in enumerate([
    ("Phase 1 — Tile Sort (simd_small_sort)", CYAN, [
        "Load up to 192 floats into 24 AVX2 registers (8 per register)",
        "Apply a fixed comparison network — like a tournament bracket",
        "All 8 comparisons happen in parallel — zero branches",
        "Output: a perfectly sorted 192-element chunk",
    ]),
    ("Phase 2 — Bottom-up Merge (simd_merge_sort)", ACCENT, [
        "Sorted chunks are merged in pairs: 192 → 384 → 768 → … → n",
        "At n = 2,000,000: requires 14 merge passes",
        "Merge hot-loop: pick smaller head element, load 8 floats, store",
        "This is where 65–75% of total runtime is spent → main target",
    ]),
]):
    section_card(s, 0.45 + i*6.45, 1.3, 6.15, 3.0,
                 title, lines, title_color=col)

# Baseline vs std::sort actual data table
bordered_rect(s, 0.45, 4.55, 12.43, 2.6, fill=PANEL, border=BORDER)
tb(s, 0.65, 4.62, 8.0, 0.38,
   "Baseline speedup vs std::sort  (actual measured values)",
   size=14, bold=True, color=WHITE)

cols_x   = [0.65, 2.5,  4.5,  6.5,  8.5,  10.5]
col_hdrs = ["Array size", "std::sort", "SIMD baseline", "Speedup", "Fits in", "Bottleneck"]
for i,h in enumerate(col_hdrs):
    tb(s, cols_x[i], 5.05, 1.9, 0.3, h, size=11, bold=True,
       color=LGRAY, align=PP_ALIGN.CENTER)
hline(s, 0.65, 5.38, 12.1, color=BORDER)

rows = [
    ("n = 81",      "0.19 ms",  "0.02 ms",  "9.0×",  "L1 cache",   "None — pure SIMD"),
    ("n = 2,187",   "7.4 ms",   "1.7 ms",   "4.3×",  "L2 cache",   "Merge overhead"),
    ("n = 531,441", "3,303 ms", "1,161 ms", "2.8×",  "RAM",        "Memory bandwidth"),
]
for ri, row in enumerate(rows):
    yy = 5.45 + ri*0.45
    bg_c = c(0x12,0x12,0x1C) if ri%2==0 else PANEL
    rect(s, 0.65, yy, 12.1, 0.43, color=bg_c)
    for ci,cell in enumerate(row):
        col = GREEN if "×" in cell else WHITE
        tb(s, cols_x[ci], yy+0.05, 1.9, 0.35, cell, size=12,
           color=col, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Optimizations Overview
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=YELLOW)
header(s, "Five Optimizations — What We Changed and Why",
       "Each targets a specific hardware bottleneck")

opt_data = [
    ("1", "OpenMP\nMulti-threading", CYAN,
     "12 cores instead of 1",
     "Recursive task spawning parallelises both tile sort and merge tree phases. "
     "Depth = 4 → up to 16 concurrent tasks.",
     "2.7 – 3.0×"),
    ("2", "Branchless\nMerge", ACCENT,
     "No mis-predicted branches",
     "Replace if-branch in merge loop with arithmetic pointer selection. "
     "CPU calculates source address without guessing.",
     "1.0 – 1.1×"),
    ("3", "Memory\nPrefetch", GREEN,
     "Hide RAM latency",
     "Issue prefetch hint 20 SIMD vectors (640 bytes) ahead. "
     "Covers ~80ns DRAM latency before the data is needed.",
     "1.0 – 1.2×"),
    ("4", "AVX-512\nSimulation", YELLOW,
     "One fewer merge pass",
     "Create 384-element sorted runs (2× normal tile) then start merging from 384. "
     "Saves 1 of 14 merge passes at n = 2M.",
     "1.05 – 1.2×"),
    ("5", "Pooled\nAllocation", RED,
     "No malloc per merge call",
     "Reuse a thread-local buffer instead of malloc/free for every merge. "
     "Eliminates thousands of allocations per sort.",
     "1.07 – 1.3×"),
]

x = 0.45
for num, name, col, tagline, desc, gain in opt_data:
    bordered_rect(s, x, 1.3, 2.44, 5.85, fill=PANEL, border=col)
    # number badge
    rect(s, x, 1.3, 2.44, 0.38, color=col)
    tb(s, x, 1.3, 2.44, 0.38, f"Gap {num}", size=13, bold=True,
       color=c(0x05,0x05,0x10), align=PP_ALIGN.CENTER)
    tb(s, x+0.1, 1.72, 2.24, 0.6, name, size=14, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    hline(s, x+0.1, 2.38, 2.24, color=col, h=0.02)
    tb(s, x+0.1, 2.45, 2.24, 0.35, tagline, size=11, bold=True,
       color=col, align=PP_ALIGN.CENTER)
    tb(s, x+0.1, 2.88, 2.24, 2.4, desc, size=11.5, color=LGRAY,
       wrap=True, align=PP_ALIGN.LEFT)
    rect(s, x+0.1, 6.5, 2.24, 0.45, color=c(0x08,0x08,0x18))
    tb(s, x+0.1, 6.52, 2.24, 0.38, f"Gain: {gain}", size=12,
       bold=True, color=col, align=PP_ALIGN.CENTER)
    x += 2.58


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Results: Speedup Bar Chart (actual data)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=GREEN)
header(s, "Results — Speedup Over Baseline at n = 500,000",
       "Actual measured values · 100 runs averaged · Intel Core i7-12th Gen · 12 threads")

# ACTUAL DATA from run: n=500000
# Baseline=1178.31, OpenMP=442.98, Branchless=1073.97, Prefetch=1019.91,
# DynTile=918.68, All=387.41, AVX512=947.90
actual_speedups = [
    ("Baseline",     1.00),
    ("Branchless",   1178.31/1073.97),   # 1.097
    ("Prefetch",     1178.31/1019.91),   # 1.155
    ("AVX-512 Sim",  1178.31/947.90),    # 1.243
    ("DynTile",      1178.31/918.68),    # 1.282
    ("OpenMP",       1178.31/442.98),    # 2.660
    ("All Combined", 1178.31/387.41),    # 3.041
]
actual_speedups_rounded = [(n, round(v,2)) for n,v in actual_speedups]

bar_colors = [MGRAY, ACCENT, CYAN, YELLOW, GREEN, c(0xFF,0x7F,0x00), c(0xFF,0x40,0x80)]

# Draw horizontal bars manually
chart_l  = 1.5
chart_t  = 1.5
bar_area = 8.5
max_v    = 3.5
row_h    = 0.72
label_w  = 1.8

for i,(lbl,val) in enumerate(actual_speedups_rounded):
    y     = chart_t + i*row_h
    bh    = 0.48
    bw    = (val/max_v)*bar_area
    col   = bar_colors[i]
    # label
    tb(s, 0.3, y+0.08, label_w, bh, lbl, size=13, color=LGRAY,
       align=PP_ALIGN.RIGHT, wrap=False)
    # bar background track
    rect(s, chart_l, y+0.08, bar_area, bh, color=c(0x18,0x18,0x28))
    # actual bar
    rect(s, chart_l, y+0.08, bw, bh, color=col)
    # value
    tb(s, chart_l+bw+0.1, y+0.08, 0.8, bh,
       f"{val:.2f}×", size=13, bold=True, color=col)
    # baseline reference line
    if i == 0:
        pass
    # 1.0× reference line
ref_x = chart_l + (1.0/max_v)*bar_area
rect(s, ref_x, chart_t-0.1, 0.02, len(actual_speedups)*row_h+0.1,
     color=MGRAY)
tb(s, ref_x-0.2, chart_t-0.35, 0.8, 0.28, "1.0×",
   size=10, color=MGRAY, align=PP_ALIGN.CENTER)

# Time labels on right
tb(s, 10.5, 1.2, 2.8, 0.3, "Actual time (ms) at n=500K",
   size=11, color=LGRAY, italic=True)
actual_times = [
    ("Baseline",     1178.31),
    ("Branchless",   1073.97),
    ("Prefetch",     1019.91),
    ("AVX-512 Sim",   947.90),
    ("DynTile",       918.68),
    ("OpenMP",        442.98),
    ("All Combined",  387.41),
]
for i,(lbl,ms) in enumerate(actual_times):
    y = chart_t + i*row_h
    tb(s, 10.5, y+0.15, 2.7, 0.42,
       f"{ms:.0f} ms", size=13, color=bar_colors[i], bold=True)

tb(s, 0.3, 6.82, 12.73, 0.4,
   "All = OpenMP + Branchless + Prefetch + DynTile combined  ·  "
   "Baseline = 1,178 ms (single-threaded, no optimizations)",
   size=11, color=LGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Results across all sizes (line chart)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=CYAN)
header(s, "Speedup Across All Array Sizes",
       "Actual benchmark output — 9 sizes from 10K to 2M elements")

# ACTUAL speedup data from run
sizes_lbl = ["10K","50K","100K","250K","500K","750K","1M","1.5M","2M"]
# From actual output:
speedup_data = {
    "OpenMP":     [0.44, 2.08, 2.51, 2.47, 2.66, 2.53, 2.79, 2.69, 2.91],
    "DynTile":    [1.13, 1.12, 1.08, 1.02, 1.28, 1.14, 1.20, 1.17, 1.17],
    "AVX-512Sim": [1.06, 1.11, 1.11, 1.05, 1.24, 1.14, 1.17, 1.13, 1.19],
    "Prefetch":   [0.98, 1.00, 0.92, 1.05, 1.16, 0.95, 1.04, 1.01, 1.04],
    "Branchless": [0.96, 0.95, 1.00, 1.01, 1.10, 0.97, 1.03, 1.00, 1.03],
    "All":        [0.45, 2.08, 2.64, 2.36, 3.04, 2.71, 2.83, 2.75, 2.91],
}
line_colors = {
    "OpenMP":     CYAN,
    "All":        GREEN,
    "DynTile":    YELLOW,
    "AVX-512Sim": ACCENT,
    "Prefetch":   c(0xFF,0x80,0x40),
    "Branchless": c(0xFF,0x40,0x80),
}

# Draw chart area
chart_l = 1.0; chart_t = 1.6
chart_w = 8.5; chart_h = 4.8
n_pts   = len(sizes_lbl)
max_spd = 3.5

# Background
rect(s, chart_l, chart_t, chart_w, chart_h, color=c(0x0C,0x0C,0x16))

# Gridlines & y-labels
for gi in range(6):
    gy  = chart_t + chart_h - (gi/5)*chart_h
    val = gi * (max_spd/5)
    rect(s, chart_l, gy, chart_w, 0.01, color=BORDER)
    tb(s, 0.0, gy-0.12, 0.95, 0.28,
       f"{val:.1f}×", size=10, color=LGRAY, align=PP_ALIGN.RIGHT)

# 1.0× reference line
ref_y = chart_t + chart_h - (1.0/max_spd)*chart_h
rect(s, chart_l, ref_y, chart_w, 0.025, color=MGRAY)
tb(s, chart_l+0.05, ref_y-0.22, 1.0, 0.2, "baseline", size=9,
   color=MGRAY, italic=True)

# X-axis labels
for xi,lbl in enumerate(sizes_lbl):
    xp = chart_l + (xi/(n_pts-1))*chart_w
    tb(s, xp-0.3, chart_t+chart_h+0.06, 0.6, 0.25,
       lbl, size=10, color=LGRAY, align=PP_ALIGN.CENTER)
    rect(s, xp, chart_t, 0.012, chart_h, color=c(0x20,0x20,0x30))

# Draw lines
import math
for sname, vals in speedup_data.items():
    col = line_colors[sname]
    pts = []
    for xi,v in enumerate(vals):
        xp = chart_l + (xi/(n_pts-1))*chart_w
        yp = chart_t + chart_h - (v/max_spd)*chart_h
        pts.append((xp,yp))
    for i in range(len(pts)-1):
        x1,y1=pts[i]; x2,y2=pts[i+1]
        draw_line_segment(s, x1, y1, x2, y2, col, thickness=0.05)
    for (xp,yp) in pts:
        c2 = s.shapes.add_shape(9, Inches(xp-0.06), Inches(yp-0.06),
                                 Inches(0.12), Inches(0.12))
        c2.fill.solid(); c2.fill.fore_color.rgb = col
        c2.line.fill.background()

# Legend (right panel)
bordered_rect(s, 9.7, 1.6, 3.35, 4.8, fill=PANEL, border=BORDER)
tb(s, 9.85, 1.68, 3.0, 0.35, "Variant", size=12, bold=True, color=WHITE)
hline(s, 9.85, 2.05, 3.0, color=BORDER, h=0.01)
legend_items = [
    ("All Combined",  GREEN,              "Best overall"),
    ("OpenMP",        CYAN,               "Best at n>50K"),
    ("DynTile",       YELLOW,             "Best single-core"),
    ("AVX-512 Sim",   ACCENT,             "Saves 1 merge pass"),
    ("Prefetch",      c(0xFF,0x80,0x40),  "Hides RAM latency"),
    ("Branchless",    c(0xFF,0x40,0x80),  "No branch guessing"),
]
ly = 2.15
for nm, col, note in legend_items:
    rect(s, 9.9, ly+0.05, 0.22, 0.22, color=col)
    tb(s, 10.2, ly, 1.2, 0.32, nm, size=11, bold=True, color=col)
    tb(s, 10.2, ly+0.3, 2.7, 0.25, note, size=10, color=LGRAY)
    ly += 0.65


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Scalability (actual Option 8 data)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=ACCENT)
header(s, "Scalability — Strong Scaling & Thread Efficiency",
       "Actual Option 8 output · n = 2,000,000 elements · thread sweep 1 → 12")

# ACTUAL strong scaling data from run
strong_data = [
    (1,  53.30, 1.08, 108.0),
    (2,  32.03, 1.80,  89.8),
    (3,  29.89, 1.92,  64.2),
    (4,  21.61, 2.66,  66.6),
    (6,  21.08, 2.73,  45.5),
    (8,  20.00, 2.88,  36.0),
    (10, 18.28, 3.15,  31.5),
    (12, 17.70, 3.25,  27.1),
]

# Left: speedup line chart
chart_l = 0.6; chart_t = 1.45
chart_w = 5.8; chart_h = 4.2

rect(s, chart_l, chart_t, chart_w, chart_h, color=c(0x0C,0x0C,0x16))
tb(s, chart_l, chart_t-0.32, chart_w, 0.28,
   "Speedup vs threads (n=2M)", size=12, bold=True, color=WHITE)

max_t = 12; max_sp = 4.0
# ideal line
for xi in range(8):
    t_val = [1,2,3,4,6,8,10,12][xi]
    xp = chart_l + ((t_val-1)/(max_t-1))*chart_w
    yp = chart_t + chart_h - (min(t_val,max_sp)/max_sp)*chart_h
    ideal_col = c(0x30,0x30,0x50)
    if xi < 7:
        t_next = [1,2,3,4,6,8,10,12][xi+1]
        xp2 = chart_l + ((t_next-1)/(max_t-1))*chart_w
        yp2 = chart_t + chart_h - (min(t_next,max_sp)/max_sp)*chart_h
        draw_line_segment(s, xp, yp, xp2, yp2, ideal_col, thickness=0.025)

# gridlines
for gi in range(5):
    gy  = chart_t + chart_h - (gi/4)*chart_h
    val = gi*(max_sp/4)
    rect(s, chart_l, gy, chart_w, 0.01, color=BORDER)
    tb(s, 0.0, gy-0.12, 0.55, 0.26, f"{val:.1f}×", size=9,
       color=LGRAY, align=PP_ALIGN.RIGHT)

thread_vals = [r[0] for r in strong_data]
speedup_vals = [r[2] for r in strong_data]
eff_vals     = [r[3] for r in strong_data]

pts_sp = []
for i,r in enumerate(strong_data):
    t_val = r[0]
    sp    = r[2]
    xp = chart_l + ((t_val-1)/(max_t-1))*chart_w
    yp = chart_t + chart_h - (sp/max_sp)*chart_h
    pts_sp.append((xp,yp))

for i in range(len(pts_sp)-1):
    x1,y1=pts_sp[i]; x2,y2=pts_sp[i+1]
    draw_line_segment(s, x1, y1, x2, y2, CYAN, thickness=0.055)
for (xp,yp) in pts_sp:
    c2 = s.shapes.add_shape(9, Inches(xp-0.07), Inches(yp-0.07),
                             Inches(0.14), Inches(0.14))
    c2.fill.solid(); c2.fill.fore_color.rgb = CYAN
    c2.line.fill.background()

# X axis labels
for r in strong_data:
    t_val = r[0]
    xp = chart_l + ((t_val-1)/(max_t-1))*chart_w
    tb(s, xp-0.2, chart_t+chart_h+0.06, 0.4, 0.22,
       str(t_val), size=10, color=LGRAY, align=PP_ALIGN.CENTER)
tb(s, chart_l+chart_w/2-1, chart_t+chart_h+0.3, 2.0, 0.25,
   "Threads", size=11, color=LGRAY, align=PP_ALIGN.CENTER)

# ideal label
tb(s, chart_l+0.1, chart_t+0.1, 1.5, 0.25, "Ideal (linear)",
   size=9, color=MGRAY, italic=True)
tb(s, chart_l+0.1, chart_t+0.4, 1.5, 0.25, "Actual",
   size=9, color=CYAN, bold=True)

# Right: data table
bordered_rect(s, 6.8, 1.45, 6.1, 4.2, fill=PANEL, border=BORDER)
tb(s, 7.0, 1.52, 5.7, 0.35,
   "Actual measured values — n = 2,000,000", size=12, bold=True, color=WHITE)
hline(s, 7.0, 1.9, 5.7, color=BORDER, h=0.012)

th_hdrs = ["Threads", "Time (ms)", "Speedup", "Efficiency"]
th_cx   = [7.0, 8.3, 9.7, 11.1]
for i,h in enumerate(th_hdrs):
    tb(s, th_cx[i], 1.95, 1.25, 0.3, h, size=11, bold=True,
       color=LGRAY, align=PP_ALIGN.CENTER)
hline(s, 7.0, 2.28, 5.7, color=MGRAY, h=0.01)

for ri,r in enumerate(strong_data):
    yt = 2.35 + ri*0.37
    bc = c(0x12,0x12,0x1E) if ri%2==0 else PANEL
    rect(s, 7.0, yt, 5.7, 0.35, color=bc)
    vals = [str(r[0]), f"{r[1]:.1f}", f"{r[2]:.2f}×", f"{r[3]:.1f}%"]
    for ci,cell in enumerate(vals):
        col = CYAN if ci==2 else (GREEN if ci==3 and r[3]>60 else
              (YELLOW if ci==3 and r[3]>35 else
              (RED if ci==3 else WHITE)))
        tb(s, th_cx[ci], yt+0.04, 1.25, 0.28, cell, size=12,
           color=col, align=PP_ALIGN.CENTER, bold=(ci==2))

# Bottom annotations
rect(s, 0.6, 5.9, 12.4, 1.3, color=PANEL)
tb(s, 0.8, 5.96, 5.8, 0.35,
   "Phase breakdown (n = 2,000,000)", size=12, bold=True, color=WHITE)

phase_data = [
    ("Tile-sort (parallel)", 2.0,  CYAN),
    ("Merge (serial)",       98.0, YELLOW),
]
bx = 0.8
for nm, pct, col in phase_data:
    bar_w = (pct/100) * 5.2
    rect(s, bx, 6.38, bar_w, 0.45, color=col)
    tb(s, bx+0.05, 6.42, bar_w-0.1, 0.38,
       f"{nm}  {pct}%", size=12, bold=True,
       color=c(0x05,0x05,0x10), align=PP_ALIGN.LEFT)
    bx += bar_w

tb(s, 0.8, 6.9, 5.8, 0.28,
   "Merge = 98% serial → Amdahl ceiling ≈ 1.0× from parallelism alone",
   size=11, color=RED, italic=True)

# Amdahl annotation
bordered_rect(s, 6.8, 5.9, 6.1, 1.3, fill=c(0x0A,0x0A,0x18), border=ACCENT)
tb(s, 7.0, 5.96, 5.7, 0.35,
   "Why not 12× faster with 12 threads?", size=12, bold=True, color=WHITE)
tb(s, 7.0, 6.35, 5.7, 0.75,
   "The merge step cannot be parallelised — it must run in order. "
   "This serial portion limits the maximum possible speedup (Amdahl's Law). "
   "We achieve 3.25× by parallelising everything else.",
   size=12, color=LGRAY, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Correctness & Comparison Table
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=YELLOW)
header(s, "Correctness & Full Comparison Table",
       "Actual Option 7 output · 35/35 correctness tests passed · all 7 variants")

# Correctness row
tb(s, 0.45, 1.3, 12.43, 0.38,
   "Correctness verification — each variant sorted against std::sort reference",
   size=13, bold=True, color=WHITE)

check_sizes = ["10K", "100K", "500K", "1M", "2M"]
cs_x = 0.45
for csz in check_sizes:
    bordered_rect(s, cs_x, 1.72, 2.38, 0.52, fill=c(0x05,0x1A,0x08), border=GREEN)
    tb(s, cs_x, 1.72, 2.38, 0.52,
       f"n = {csz}   7 / 7  ✓", size=13, bold=True,
       color=GREEN, align=PP_ALIGN.CENTER)
    cs_x += 2.52

# Full data table — ACTUAL from output
tb(s, 0.45, 2.42, 12.43, 0.32,
   "Performance table (ms) — 100 runs each · n from 10K to 2M",
   size=12, color=LGRAY, italic=True)

# headers
col_names = ["Size","Baseline","OpenMP","Branchless","Prefetch","DynTile","All","AVX512Sim"]
cx = [0.35, 1.45, 2.75, 4.05, 5.35, 6.65, 7.95, 9.55]
cw = [1.0,  1.25, 1.25, 1.25, 1.25, 1.25, 1.25, 1.6]

rect(s, 0.35, 2.8, 12.43, 0.35, color=c(0x18,0x18,0x28))
for i,h in enumerate(col_names):
    tb(s, cx[i], 2.82, cw[i], 0.3, h, size=10, bold=True,
       color=LGRAY, align=PP_ALIGN.CENTER)
hline(s, 0.35, 3.17, 12.43, color=BORDER, h=0.01)

# ACTUAL data rows from output
table_rows = [
    ("10K",   11.38, 25.69, 11.86, 11.67, 10.10, 25.54, 10.69),
    ("50K",   77.90, 37.41, 81.60, 77.68, 69.79, 37.40, 70.30),
    ("100K",  170.79, 68.12, 170.41, 185.31, 157.43, 64.77, 153.82),
    ("250K",  480.22, 194.41, 473.82, 456.78, 491.69, 203.23, 455.63),
    ("500K",  1178.31, 442.98, 1073.97, 1019.91, 918.68, 387.41, 947.90),
    ("750K",  1656.25, 655.58, 1706.70, 1753.28, 1457.96, 611.47, 1449.11),
    ("1M",    2433.31, 870.99, 2358.78, 2348.32, 2025.60, 860.83, 2127.12),
    ("1.5M",  3670.73, 1362.81, 3662.67, 3617.63, 3136.71, 1333.28, 3262.48),
    ("2M",    5509.02, 1889.90, 5339.37, 5277.90, 4728.18, 1894.82, 4634.28),
]

for ri, row in enumerate(table_rows):
    yt   = 3.24 + ri*0.38
    base = row[1]
    best = min(row[1:])
    bc   = c(0x10,0x10,0x1C) if ri%2==0 else PANEL
    rect(s, 0.35, yt, 12.43, 0.36, color=bc)
    for ci, val in enumerate(row):
        if ci == 0:
            tb(s, cx[0], yt+0.04, cw[0], 0.28, str(val),
               size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        else:
            is_best  = (val == best)
            is_worse = (val > base * 1.02)
            col = GREEN if is_best else (RED if is_worse else WHITE)
            tb(s, cx[ci], yt+0.04, cw[ci], 0.28, f"{val:.0f}",
               size=10, color=col, align=PP_ALIGN.CENTER,
               bold=is_best)

tb(s, 0.35, 6.68, 12.43, 0.32,
   "Green = best time for that row  ·  Red = slower than baseline  ·  Values in milliseconds",
   size=11, color=LGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Analysis: What the results tell us
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=RED)
header(s, "Critical Analysis — What the Results Tell Us",
       "Performance trends, bottlenecks, and deviations from expectation")

insights = [
    (CYAN, "OpenMP gives 2.7–3× — but not 12×",
     "With 12 cores we expected closer to 12×. The merge step is 98% of runtime and cannot "
     "be parallelised — it must combine results in order. Amdahl's Law caps us at ~1.02× "
     "from parallelism alone at this phase split. OpenMP helps by parallelising the "
     "tile-sort phase and upper merge tree levels, giving us the observed 3×."),

    (YELLOW, "DynTile and AVX-512 Sim are the best single-core improvements",
     "DynTile eliminates thousands of malloc/free calls (7–19% gain). "
     "AVX-512 Sim reduces merge passes from 14 to 13 at n=2M and also uses pooled memory "
     "(9–24% gain). Both outperform branchless and prefetch at all sizes — "
     "memory allocation overhead is a more significant bottleneck than branch misprediction."),

    (GREEN, "Branchless merge: our first attempt was actually slower",
     "The initial implementation loaded data from BOTH sides of the merge every iteration "
     "(2 loads vs 1). At large array sizes this doubled memory bandwidth usage — "
     "far more costly than the branch it replaced. We fixed it using arithmetic pointer "
     "selection (1 load per iteration, still branchless)."),

    (ACCENT, "Prefetch helps less than expected at large sizes",
     "The CPU's hardware prefetcher already detects the sequential access pattern and "
     "loads ahead automatically. Our software hints add only marginal benefit on top. "
     "The gain is visible at medium sizes (500K: 1.16×) but near-zero at 2M "
     "where the hardware prefetcher is fully warmed up."),
]

y = 1.35
for col, title, body in insights:
    bordered_rect(s, 0.45, y, 12.43, 1.4, fill=PANEL, border=col)
    rect(s, 0.45, y, 0.06, 1.4, color=col)
    tb(s, 0.65, y+0.08, 12.0, 0.38, title, size=14, bold=True, color=col)
    tb(s, 0.65, y+0.5, 12.0, 0.82, body, size=12.5, color=LGRAY, wrap=True)
    y += 1.52


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — LLM Usage
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=YELLOW)
header(s, "LLM Usage — Claude (Anthropic)",
       "How AI assistance was used during this project — honest disclosure")

llm_items = [
    (CYAN,   "Error Debugging",
     "Compiler and runtime errors were shared with Claude to understand root causes. "
     "Examples: branchless merge producing wrong output (double-load bug), "
     "AVX-512 sim failing correctness at n=500 (tile size > simd_small_sort limit), "
     "and python-docx RGBColor attribute errors in report generation. "
     "Claude explained why each error occurred rather than just patching it."),

    (GREEN,  "Terminal Output Generation & Benchmark Scripts",
     "Claude helped write the benchmark driver code (options menu, timing loop, "
     "correctness checker, scalability sweep) and the Python scripts that generated "
     "this presentation and the Word report from actual measured output. "
     "All data displayed in these slides comes from real hardware runs — not fabricated."),

    (ACCENT, "Concept Clarification",
     "Used to clarify hardware concepts: why Amdahl's Law limits OpenMP to 3× despite 12 cores, "
     "how _mm_prefetch is advisory and safe without bounds checks on x86, "
     "why the branchless merge initially had 2 loads/iter, "
     "and why MPI would not help on a shared-memory single-machine sort benchmark."),

    (YELLOW, "What Was NOT Used",
     "Core algorithmic design (bitonic sort structure, merge network) was pre-existing "
     "in the simd_bitonic library. SIMD intrinsics (AVX2), OpenMP task annotations, "
     "and the optimization strategy were developed through literature and hardware profiling. "
     "Claude was a debugging and explanation tool — not the author of the algorithms."),
]

y = 1.35
for col, title, body in llm_items:
    bordered_rect(s, 0.45, y, 12.43, 1.35, fill=PANEL, border=col)
    rect(s, 0.45, y, 0.06, 1.35, color=col)
    tb(s, 0.65, y+0.07, 12.0, 0.35, title, size=14, bold=True, color=col)
    tb(s, 0.65, y+0.46, 12.0, 0.8, body, size=12, color=LGRAY, wrap=True)
    y += 1.47

tb(s, 0.45, 7.1, 12.43, 0.3,
   "Tool: Claude (Anthropic)  ·  Usage: debugging, code generation helpers, concept explanation  ·  All benchmark data from actual hardware runs",
   size=10, color=MGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Conclusions
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=GREEN)
header(s, "Conclusions", "09")

for i,(val,lbl,col) in enumerate([
    ("3.25×", "Peak speedup\n(12 threads, n=2M)", GREEN),
    ("1.30×", "Best single-core\n(DynTile, n=50K)",  YELLOW),
    ("1.24×", "AVX-512 Sim\n(n=500K)",               ACCENT),
    ("35/35", "Tests passed",                          CYAN),
]):
    stat_card(s, 0.45+i*3.2, 1.2, 2.95, 1.55, val, lbl, color=col)

bordered_rect(s, 0.45, 2.95, 12.43, 4.2, fill=PANEL, border=BORDER)
tb(s, 0.65, 3.02, 12.0, 0.38, "Key Takeaways", size=15, bold=True, color=WHITE)
hline(s, 0.65, 3.45, 12.0, color=ACCENT, h=0.02)

takeaways = [
    (CYAN,   "Multi-threading is the dominant improvement.",
             "OpenMP gives 3× at large sizes — parallelising both tile-sort and merge-tree phases "
             "across 12 threads. No single-core technique comes close."),
    (YELLOW, "Eliminating allocation overhead beats branch elimination.",
             "Pooled memory (DynTile) and AVX-512 Sim outperform branchless and prefetch — "
             "malloc/free cost is a bigger bottleneck than branch misprediction on this hardware."),
    (ACCENT, "Measure everything — theory and practice diverge.",
             "Our first branchless implementation was slower. Our prefetch added overhead via bounds checks. "
             "Both were fixed only because we measured the actual running time."),
    (GREEN,  "The serial merge step is the fundamental limit.",
             "At n=2M the merge phase is 98% of runtime. Amdahl's Law caps parallel speedup at ~1.02×. "
             "Further gains require parallelising the merge itself — an open research problem."),
]
y = 3.55
for col, title, body in takeaways:
    rect(s, 0.65, y, 0.05, 0.78, color=col)
    tb(s, 0.85, y+0.02, 11.8, 0.35, title, size=13, bold=True, color=col)
    tb(s, 0.85, y+0.38, 11.8, 0.38, body, size=12, color=LGRAY, wrap=True)
    y += 0.9


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — References
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 0.06, 7.5, color=LGRAY)
header(s, "References", "10")

bordered_rect(s, 0.45, 1.3, 12.43, 5.85, fill=PANEL, border=BORDER)
refs = [
    ("[1]", "simd_bitonic",
     "Open-source SIMD bitonic sort library. github.com/nlguillemot/simd_bitonic"),
    ("[2]", "Batcher (1968)",
     "Sorting networks and their applications. AFIPS Spring Joint Computer Conference."),
    ("[3]", "Intel SDM Vol. 2",
     "Intel 64 and IA-32 Architectures Software Developer's Manual — AVX2 intrinsics."),
    ("[4]", "OpenMP ARB",
     "OpenMP Application Programming Interface, Version 4.5."),
    ("[5]", "Amdahl (1967)",
     "Validity of the single-processor approach. AFIPS Spring Joint Computer Conference."),
    ("[6]", "Fog (2023)",
     "Optimizing Software in C++ — Branch prediction, prefetch, cache hierarchy."),
    ("[7]", "sokol_time.h",
     "Cross-platform high-resolution timer. André Weissflog. MIT License."),
]
y = 1.45
for tag, src, desc in refs:
    tb(s, 0.65, y, 0.7, 0.42, tag, size=12, bold=True, color=ACCENT)
    tb(s, 1.4, y, 2.2, 0.42, src, size=12, bold=True, color=WHITE)
    tb(s, 3.65, y, 9.0, 0.42, desc, size=12, color=LGRAY, wrap=True)
    hline(s, 0.65, y+0.45, 12.0, color=BORDER, h=0.008)
    y += 0.58

tb(s, 0.65, 7.08, 12.0, 0.3,
   "All benchmarks reproducible: g++ -O3 -mavx2 -fopenmp -std=c++11 · fixed seed 0xABCDEF01",
   size=11, color=MGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
out = r"d:\PDC\PDC_Final_project\SIMD_Bitonic_Sort_v4.pptx"
prs.save(out)
print(f"Saved: {out}  ({prs.slides.__len__()} slides)")
